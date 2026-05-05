"""Generate the English CV PPTX template (A4 portrait) for Bonfire."""

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# A4 dimensions in EMU
A4_WIDTH = Emu(210 * 36000)   # 210 mm
A4_HEIGHT = Emu(297 * 36000)  # 297 mm

# Margins
LEFT_MARGIN = Mm(20)
RIGHT_MARGIN = Mm(20)
TOP_MARGIN = Mm(15)

# Colors
DARK_BLUE = RGBColor(0x1A, 0x2B, 0x4A)  # Dark navy blue for headers
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)

# Photo
PHOTO_PATH = "templates/en/J.Aponte_Foto.jpg"
PHOTO_WIDTH = Mm(45)
PHOTO_HEIGHT = Mm(55)


def _set_slide_size(prs):
    """Set custom A4 portrait slide size."""
    prs.slide_width = A4_WIDTH
    prs.slide_height = A4_HEIGHT


def _add_text_box(slide, left, top, width, height, text, font_size=11,
                  bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a text box with specified formatting."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return shape


def _add_horizontal_line(slide, left, top, width, color=DARK_BLUE, thickness=Pt(1.5)):
    """Add a horizontal line shape."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.line.width = Pt(0)
    return line


def _add_section_header(slide, left, top, width, title):
    """Add a section header with a horizontal line above it."""
    line_top = top
    text_top = top + Mm(2)
    _add_horizontal_line(slide, left, line_top, width)
    _add_text_box(
        slide, left, text_top, width, Mm(8),
        text=title,
        font_size=14,
        bold=True,
        color=DARK_BLUE,
    )
    return text_top + Mm(8)


def build_template():
    """Build the A4 portrait CV template with placeholders."""
    prs = Presentation()
    _set_slide_size(prs)

    # Use a blank layout
    blank_layout = prs.slide_layouts[6]

    # ================================================================
    # PAGE 1
    # ================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    current_top = TOP_MARGIN

    # --- Photo (top left) ---
    photo_left = LEFT_MARGIN
    photo_top = current_top
    slide1.shapes.add_picture(
        PHOTO_PATH, photo_left, photo_top, width=PHOTO_WIDTH, height=PHOTO_HEIGHT
    )

    # --- Name (top right of photo) ---
    name_left = photo_left + PHOTO_WIDTH + Mm(8)
    name_top = photo_top + Mm(2)
    _add_text_box(
        slide1, name_left, name_top, Mm(120), Mm(12),
        text="[personal_info_name]",
        font_size=28,
        bold=True,
        color=DARK_BLUE,
    )

    # --- Title / Role ---
    title_top = name_top + Mm(10)
    _add_text_box(
        slide1, name_left, title_top, Mm(120), Mm(8),
        text="[personal_info_title]",
        font_size=14,
        bold=False,
        color=DARK_BLUE,
    )

    # --- Contact Info ---
    contact_top = title_top + Mm(10)
    contact_lines = [
        "Email:      [personal_info_email]",
        "Phone:      [personal_info_phone]",
        "Address:    [personal_info_address]",
        "LinkedIn:   [personal_info_linkedin]",
    ]
    for i, line in enumerate(contact_lines):
        _add_text_box(
            slide1, name_left, contact_top + Mm(5 * i), Mm(120), Mm(5),
            text=line,
            font_size=10,
            color=BLACK,
        )

    current_top = photo_top + PHOTO_HEIGHT + Mm(5)

    # --- Horizontal line under header ---
    _add_horizontal_line(
        slide1, LEFT_MARGIN, current_top, A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN
    )
    current_top += Mm(4)

    # --- PROFESSIONAL SUMMARY ---
    current_top = _add_section_header(
        slide1, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "PROFESSIONAL SUMMARY"
    )
    _add_text_box(
        slide1, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(25),
        text="[professional_summary]",
        font_size=11,
        color=BLACK,
    )
    current_top += Mm(28)

    # --- WORK EXPERIENCE ---
    current_top = _add_section_header(
        slide1, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "WORK EXPERIENCE"
    )

    # Job entry 1: Self-employed / Further education
    job1_date_top = current_top
    _add_text_box(
        slide1, LEFT_MARGIN, job1_date_top, Mm(35), Mm(20),
        text="[work_exp_1_date]\n[work_exp_1_location]",
        font_size=10,
        color=GRAY,
        align=PP_ALIGN.LEFT,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job1_date_top,
        Mm(130), Mm(8),
        text="[work_exp_1_title]",
        font_size=11,
        bold=True,
        color=BLACK,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job1_date_top + Mm(6),
        Mm(130), Mm(30),
        text="[work_descriptions:work_exp_1]",
        font_size=10,
        color=BLACK,
    )
    current_top = job1_date_top + Mm(35)

    # Job entry 2: Master thesis
    job2_date_top = current_top
    _add_text_box(
        slide1, LEFT_MARGIN, job2_date_top, Mm(35), Mm(20),
        text="[work_exp_2_date]\n[work_exp_2_location]",
        font_size=10,
        color=GRAY,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job2_date_top,
        Mm(130), Mm(8),
        text="[work_exp_2_title]",
        font_size=11,
        bold=True,
        color=BLACK,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job2_date_top + Mm(6),
        Mm(130), Mm(30),
        text="[work_descriptions:work_exp_2]",
        font_size=10,
        color=BLACK,
    )
    current_top = job2_date_top + Mm(35)

    # Job entry 3: Internship
    job3_date_top = current_top
    _add_text_box(
        slide1, LEFT_MARGIN, job3_date_top, Mm(35), Mm(20),
        text="[work_exp_3_date]\n[work_exp_3_location]",
        font_size=10,
        color=GRAY,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job3_date_top,
        Mm(130), Mm(8),
        text="[work_exp_3_title]",
        font_size=11,
        bold=True,
        color=BLACK,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job3_date_top + Mm(6),
        Mm(130), Mm(30),
        text="[work_descriptions:work_exp_3]",
        font_size=10,
        color=BLACK,
    )
    current_top = job3_date_top + Mm(35)

    # Job entry 4: Student assistant
    job4_date_top = current_top
    _add_text_box(
        slide1, LEFT_MARGIN, job4_date_top, Mm(35), Mm(20),
        text="[work_exp_4_date]\n[work_exp_4_location]",
        font_size=10,
        color=GRAY,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job4_date_top,
        Mm(130), Mm(8),
        text="[work_exp_4_title]",
        font_size=11,
        bold=True,
        color=BLACK,
    )
    _add_text_box(
        slide1, LEFT_MARGIN + Mm(40), job4_date_top + Mm(6),
        Mm(130), Mm(30),
        text="[work_descriptions:work_exp_4]",
        font_size=10,
        color=BLACK,
    )
    current_top = job4_date_top + Mm(30)

    # --- Page number ---
    _add_text_box(
        slide1, LEFT_MARGIN, A4_HEIGHT - Mm(12),
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(5),
        text="Page 1 of 2",
        font_size=9,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )

    # ================================================================
    # PAGE 2
    # ================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    current_top = TOP_MARGIN

    # --- EDUCATION ---
    current_top = _add_section_header(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "EDUCATION"
    )
    _add_text_box(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(35),
        text="[education]",
        font_size=11,
        color=BLACK,
    )
    current_top += Mm(38)

    # --- RELEVANT SUBJECTS ---
    current_top = _add_section_header(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "RELEVANT SUBJECTS"
    )
    _add_text_box(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(20),
        text="[relevant_subjects]",
        font_size=11,
        color=BLACK,
    )
    current_top += Mm(23)

    # --- SKILLS ---
    current_top = _add_section_header(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "SKILLS"
    )
    _add_text_box(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(35),
        text="[skills]",
        font_size=11,
        color=BLACK,
    )
    current_top += Mm(38)

    # --- PROJECTS ---
    current_top = _add_section_header(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "PROJECTS"
    )
    _add_text_box(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(30),
        text="[selected_projects]",
        font_size=11,
        color=BLACK,
    )
    current_top += Mm(33)

    # --- CERTIFICATIONS (optional static zone) ---
    current_top = _add_section_header(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN,
        "CERTIFICATIONS & TRAINING"
    )
    _add_text_box(
        slide2, LEFT_MARGIN, current_top,
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(20),
        text="[certifications]",
        font_size=11,
        color=BLACK,
    )

    # --- Page number ---
    _add_text_box(
        slide2, LEFT_MARGIN, A4_HEIGHT - Mm(12),
        A4_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Mm(5),
        text="Page 2 of 2",
        font_size=9,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )

    # Save
    output_path = "templates/en/cv_template.pptx"
    prs.save(output_path)
    print(f"Template saved to: {output_path}")


if __name__ == "__main__":
    build_template()
