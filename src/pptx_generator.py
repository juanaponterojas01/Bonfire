"""Render CVs from PPTX templates using placeholder replacement across all shapes.

This module fills a German `.pptx` CV template by replacing `[placeholder]`
strings with actual content. It handles the common PPTX issue where text is
split across multiple XML `<a:r>` runs inside shapes and text frames.
"""

from pathlib import Path
from pptx import Presentation

from src.models import CVDynamicZones, JobDescription, UserProfile


def _replace_text_in_shape(shape, old_text: str, new_text: str) -> None:
    """Replace *old_text* with *new_text* in a single shape's text frame.

    Handles text split across multiple XML runs within a paragraph.

    Args:
        shape: A python-pptx shape object (may have nested grouped shapes).
        old_text: The placeholder string to search for.
        new_text: The value to replace the placeholder with.
    """
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if old_text not in full_text:
                continue
            replaced = full_text.replace(old_text, new_text)
            if paragraph.runs:
                for run in paragraph.runs:
                    run.text = ""
                paragraph.runs[0].text = replaced

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _replace_text_in_shape(sub_shape, old_text, new_text)


def _replace_text_in_pptx(prs: Presentation, old_text: str, new_text: str) -> None:
    """Iterate through all slides and shapes, replacing *old_text* with *new_text*.

    Args:
        prs: The python-pptx Presentation object to modify in place.
        old_text: The placeholder string to search for.
        new_text: The value to replace the placeholder with.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_text_in_shape(shape, old_text, new_text)



def _format_relevant_subjects(dynamic_zones: CVDynamicZones) -> str:
    """Format relevant academic subjects as a bulleted list.

    Args:
        dynamic_zones: The AI-generated dynamic CV content.

    Returns:
        A string with each subject prefixed by a bullet point.
    """
    return "\n".join(f"• {subject}" for subject in dynamic_zones.relevant_subjects)


def _build_cv_context(
    dynamic_zones: CVDynamicZones,
) -> dict[str, str]:
    """Build a context dictionary mapping placeholders to formatted values.

    Args:
        dynamic_zones: The AI-generated dynamic CV content.

    Returns:
        A dictionary where keys are placeholder strings and values are
        the formatted text to replace them with.
    """
    context: dict[str, str] = {
        "[professional_summary]": dynamic_zones.professional_summary,
        "[bachelor_subjects]": dynamic_zones.bachelor_subjects,
        "[master_subjects]": dynamic_zones.master_subjects,
    }
    return context


def render_cv(
    template_path: str | Path,
    output_path: str | Path,
    dynamic_zones: CVDynamicZones,
    language: str,
) -> str:
    """Render a CV from a PPTX template by replacing placeholders.

    Args:
        template_path: Path to the .pptx template file.
        output_path: Path where the rendered CV will be saved.
        dynamic_zones: The AI-generated dynamic CV content.
        language: Language code for date formatting.

    Returns:
        The output path as a string.
    """
    prs = Presentation(str(template_path))
    context = _build_cv_context(dynamic_zones)

    for placeholder, value in context.items():
        _replace_text_in_pptx(prs, placeholder, value)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)
